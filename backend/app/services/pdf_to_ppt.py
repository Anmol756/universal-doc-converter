"""
PDF → PowerPoint conversion service.
Uses PyMuPDF to render each PDF page as an image,
then python-pptx to create a PPTX with each page as a slide.
"""

import logging
from pathlib import Path

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Emu

logger = logging.getLogger(__name__)


def convert_pdf_to_ppt(input_path: Path, output_path: Path) -> Path:
    """
    Convert a PDF file to PowerPoint (PPTX) format.

    Each PDF page is rendered as a high-resolution image and placed
    on its own slide in the resulting presentation.

    Args:
        input_path: Path to the source PDF file
        output_path: Path where the PPTX file will be saved

    Returns:
        Path to the converted PPTX file
    """
    if not input_path.exists():
        raise ValueError(f"Input file not found: {input_path}")

    if input_path.suffix.lower() != ".pdf":
        raise ValueError(f"Expected PDF file, got: {input_path.suffix}")

    output_path = output_path.with_suffix(".pptx")

    logger.info(f"Converting PDF → PPT: {input_path.name} → {output_path.name}")

    try:
        pdf_doc = fitz.open(str(input_path))
        prs = Presentation()

        for page_num in range(len(pdf_doc)):
            page = pdf_doc[page_num]

            # Render page as image at 1.5x resolution (good quality + faster)
            mat = fitz.Matrix(1.5, 1.5)
            pix = page.get_pixmap(matrix=mat)
            img_bytes = pix.tobytes("png")

            # Calculate slide dimensions to match PDF page aspect ratio
            page_width = page.rect.width
            page_height = page.rect.height
            aspect = page_height / page_width

            # Use standard slide width, scale height to match
            slide_width = Inches(10)
            slide_height = Emu(int(slide_width * aspect))

            prs.slide_width = slide_width
            prs.slide_height = slide_height

            # Add a blank slide
            blank_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_layout)

            # Save temp image for insertion
            temp_img = output_path.parent / f"_temp_page_{page_num}.png"
            temp_img.write_bytes(img_bytes)

            # Add image to slide, covering the entire slide
            slide.shapes.add_picture(
                str(temp_img), Emu(0), Emu(0),
                width=slide_width, height=slide_height
            )

            # Clean up temp image
            temp_img.unlink(missing_ok=True)

        pdf_doc.close()
        prs.save(str(output_path))

    except Exception as e:
        logger.error(f"PDF to PPT conversion failed: {e}")
        raise RuntimeError(f"Conversion failed: {str(e)}") from e

    if not output_path.exists():
        raise RuntimeError("Conversion produced no output file")

    logger.info(f"Conversion complete: {output_path.name} ({output_path.stat().st_size} bytes)")
    return output_path
