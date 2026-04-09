"""
PowerPoint → PDF conversion service.
Uses python-pptx to read slides and reportlab to generate a PDF.
Each slide's text content is extracted and rendered as a PDF page.
"""

import logging
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches
from reportlab.lib.pagesizes import landscape, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.enums import TA_LEFT

logger = logging.getLogger(__name__)


def convert_ppt_to_pdf(input_path: Path, output_path: Path) -> Path:
    """
    Convert a PowerPoint (PPTX) file to PDF format.

    Extracts text content from each slide and renders it
    as a landscape PDF document.

    Args:
        input_path: Path to the source PPTX file
        output_path: Path where the PDF file will be saved

    Returns:
        Path to the converted PDF file
    """
    if not input_path.exists():
        raise ValueError(f"Input file not found: {input_path}")

    if input_path.suffix.lower() not in (".pptx", ".ppt"):
        raise ValueError(f"Expected PowerPoint file, got: {input_path.suffix}")

    output_path = output_path.with_suffix(".pdf")

    logger.info(f"Converting PPT → PDF: {input_path.name} → {output_path.name}")

    try:
        prs = Presentation(str(input_path))

        # Create landscape PDF to match slide orientation
        pdf = SimpleDocTemplate(
            str(output_path),
            pagesize=landscape(A4),
            rightMargin=72,
            leftMargin=72,
            topMargin=72,
            bottomMargin=72,
        )

        styles = getSampleStyleSheet()
        style_title = ParagraphStyle(
            "SlideTitle",
            parent=styles["Heading1"],
            fontSize=20,
            leading=28,
            spaceAfter=16,
            fontName="Helvetica-Bold",
        )
        style_body = ParagraphStyle(
            "SlideBody",
            parent=styles["Normal"],
            fontSize=12,
            leading=18,
            spaceAfter=8,
            fontName="Helvetica",
        )
        style_slide_num = ParagraphStyle(
            "SlideNum",
            parent=styles["Normal"],
            fontSize=9,
            textColor="#888888",
            spaceAfter=4,
            fontName="Helvetica",
        )

        story = []

        for slide_idx, slide in enumerate(prs.slides):
            # Slide number header
            story.append(Paragraph(f"Slide {slide_idx + 1}", style_slide_num))

            has_content = False

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if not text:
                        continue

                    has_content = True
                    # Escape XML special chars
                    text = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

                    # Use title style for the first text block on each slide
                    if shape == slide.shapes[0] if slide.shapes else None:
                        story.append(Paragraph(text, style_title))
                    else:
                        story.append(Paragraph(text, style_body))

            if not has_content:
                story.append(Paragraph("<i>(No text content on this slide)</i>", style_body))

            # Page break between slides
            if slide_idx < len(prs.slides) - 1:
                story.append(PageBreak())

        if not story:
            story.append(Paragraph("(Empty presentation)", style_body))

        pdf.build(story)

    except Exception as e:
        logger.error(f"PPT to PDF conversion failed: {e}")
        raise RuntimeError(f"Conversion failed: {str(e)}") from e

    if not output_path.exists():
        raise RuntimeError("Conversion produced no output file")

    logger.info(f"Conversion complete: {output_path.name} ({output_path.stat().st_size} bytes)")
    return output_path
